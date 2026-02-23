"""
core/extract.py — Text extraction and chunk pipeline.

Two APIs
--------
extract_chunks(path) -> list[tuple[str, str]]
    Low-level. Returns [(ref, text), ...]. No DB access. Pure extraction.
    Stable refs (never change across re-runs):
      PDF:    p{n}        (1-indexed page; split to p{n}a/b/c if page > MAX_PDF_CHARS)
      PPTX:   s{n:02d}   (1-indexed slide)
      DOCX:   sec{n:02d} (split on Heading 1/2/3; single sec01 if no headings found)
      TXT/MD: sec{n:02d} (split on blank lines or # headings; max MAX_TXT_CHARS per section)
    Final safety pass: any chunk > MAX_TOKENS gets split with a/b/c suffix.

run_extract(conn, file_id, path, lane, cfg_obj) -> dict
    Full pipeline step:
      1. Calls extract_chunks() for TEXT_EXTRACTABLE files.
      2. Runs infer_project() + infer_typology() per file.
      3. Persists chunks (idempotent — skips unchanged, updates changed).
      4. Updates files row: project_id, typology, status=EXTRACTED.
      5. Upserts into projects table.
      6. Returns {"new": N, "updated": N, "unchanged": N, "skipped": N, "failed": N}.
    METADATA_ONLY files: creates a single "meta" chunk with a path-based surrogate.
    Never raises — wraps body in try/except.

Legacy API (backward compat)
-----------------------------
    extract(path) -> (text, surrogate)
    word_count(text) -> int
"""

from __future__ import annotations

import hashlib
import json
import logging
import re
import sqlite3
from pathlib import Path
from typing import Any

from config import cfg as _module_cfg, Config
from core.db import (
    upsert_chunk as _upsert_chunk,
    upsert_project,
    set_file_status,
    log_event,
)
from core.infer import infer_project, infer_typology

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

MAX_PDF_CHARS = 6_000    # chars per PDF page before splitting with a/b/c suffix
MAX_TXT_CHARS = 3_000    # chars per TXT/MD section before further splitting
MAX_TOKENS    = 4_000    # token safety limit for all chunk types


# ---------------------------------------------------------------------------
# Internal helpers
# ---------------------------------------------------------------------------

def _token_estimate(text: str) -> int:
    """Approximate token count: words × 4/3."""
    return len(text.split()) * 4 // 3


def _split_text(text: str, max_chars: int) -> list[str]:
    """
    Split text into parts no larger than max_chars, breaking at whitespace
    where possible. Returns at least one part even if text is empty.
    """
    text = text.strip()
    if not text:
        return [""]
    if len(text) <= max_chars:
        return [text]
    parts: list[str] = []
    while text:
        if len(text) <= max_chars:
            parts.append(text)
            break
        cut = text.rfind(" ", 0, max_chars)
        if cut <= 0:
            cut = max_chars
        parts.append(text[:cut].rstrip())
        text = text[cut:].lstrip()
    return parts if parts else [text]


def _apply_suffix(ref: str, parts: list[str]) -> list[tuple[str, str]]:
    """Return [(ref, text)] for single part, or [(refa, …), (refb, …)] for multiple."""
    if len(parts) == 1:
        return [(ref, parts[0])]
    return [(ref + chr(ord("a") + i), p) for i, p in enumerate(parts)]


def _token_safe(pairs: list[tuple[str, str]]) -> list[tuple[str, str]]:
    """
    Final safety pass: split any chunk exceeding MAX_TOKENS.
    Appends a/b/c suffix to refs that need splitting.
    """
    result: list[tuple[str, str]] = []
    max_chars = MAX_TOKENS * 5  # conservative: ~5 chars/token average
    for ref, text in pairs:
        if _token_estimate(text) <= MAX_TOKENS:
            result.append((ref, text))
        else:
            parts = _split_text(text, max_chars)
            result.extend(_apply_suffix(ref, parts))
    return result


def _make_chunk_id(file_id: str, ref: str) -> str:
    return hashlib.sha256(f"{file_id}::{ref}".encode()).hexdigest()


def _content_hash(text: str) -> str:
    return hashlib.sha256(text.encode()).hexdigest()


# ---------------------------------------------------------------------------
# extract_chunks — low-level, per-format
# ---------------------------------------------------------------------------

def extract_chunks(path: Path) -> list[tuple[str, str]]:
    """
    Extract text chunks from a supported file.
    Returns [(ref, text), ...] — stable refs, never raises.
    Returns [] for unsupported or unreadable extensions.
    """
    ext = path.suffix.lower()
    try:
        if ext == ".pdf":
            pairs = _chunks_pdf(path)
        elif ext == ".docx":
            pairs = _chunks_docx(path)
        elif ext == ".doc":
            pairs = _chunks_doc(path)
        elif ext == ".pptx":
            pairs = _chunks_pptx(path)
        elif ext == ".ppt":
            pairs = _chunks_ppt(path)
        elif ext == ".xlsx":
            pairs = _chunks_xlsx(path)
        elif ext == ".xls":
            pairs = _chunks_xls(path)
        elif ext in (".txt", ".md"):
            pairs = _chunks_txt(path)
        else:
            return []
        return _token_safe(pairs)
    except Exception as e:
        logger.warning("extract_chunks failed for %s: %s", path.name, e)
        return []


def _chunks_pdf(path: Path) -> list[tuple[str, str]]:
    import pypdf  # type: ignore

    reader = pypdf.PdfReader(str(path), strict=False)
    pairs: list[tuple[str, str]] = []
    for page_num, page in enumerate(reader.pages, start=1):
        try:
            text = (page.extract_text() or "").strip()
        except Exception:
            text = ""
        ref = f"p{page_num}"
        if len(text) > MAX_PDF_CHARS:
            parts = _split_text(text, MAX_PDF_CHARS)
            pairs.extend(_apply_suffix(ref, parts))
        else:
            pairs.append((ref, text))
    return pairs


def _chunks_docx(path: Path) -> list[tuple[str, str]]:
    from docx import Document  # type: ignore

    doc = Document(str(path))
    heading_styles = {"heading 1", "heading 2", "heading 3"}

    sections: list[list[str]] = []
    current: list[str] = []

    for para in doc.paragraphs:
        style_name = (para.style.name or "").lower()
        if style_name in heading_styles:
            if current:
                sections.append(current)
            current = [para.text] if para.text.strip() else []
        else:
            if para.text.strip():
                current.append(para.text)

    if current:
        sections.append(current)

    if not sections:
        # No headings — one big chunk
        all_text = "\n".join(
            p.text for p in doc.paragraphs if p.text.strip()
        )
        return [("sec01", all_text.strip())]

    pairs: list[tuple[str, str]] = []
    for i, lines in enumerate(sections, start=1):
        ref = f"sec{i:02d}"
        text = "\n".join(lines).strip()
        pairs.append((ref, text))
    return pairs


def _chunks_pptx(path: Path) -> list[tuple[str, str]]:
    from pptx import Presentation  # type: ignore

    prs = Presentation(str(path))
    pairs: list[tuple[str, str]] = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        ref = f"s{slide_num:02d}"
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        parts.append(line)
        text = "\n".join(parts).strip()
        pairs.append((ref, text))
    return pairs


def _chunks_doc(path: Path) -> list[tuple[str, str]]:
    """
    Word 97-2003 (.doc) — try python-docx first (works if the file is
    OOXML mislabelled as .doc); otherwise return [] for graceful fallback.
    """
    try:
        from docx import Document  # type: ignore
        doc = Document(str(path))
        heading_styles = {"heading 1", "heading 2", "heading 3"}
        sections: list[list[str]] = []
        current: list[str] = []
        for para in doc.paragraphs:
            style_name = (para.style.name or "").lower()
            if style_name in heading_styles:
                if current:
                    sections.append(current)
                current = [para.text] if para.text.strip() else []
            elif para.text.strip():
                current.append(para.text)
        if current:
            sections.append(current)
        if not sections:
            all_text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            return [("sec01", all_text.strip())] if all_text.strip() else []
        return [(f"sec{i:02d}", "\n".join(lines).strip())
                for i, lines in enumerate(sections, start=1)]
    except Exception:
        return []   # True binary .doc — no pure-Python parser available


def _chunks_ppt(path: Path) -> list[tuple[str, str]]:
    """PowerPoint 97-2003 (.ppt) — no pure-Python parser; return empty."""
    return []


def _chunks_xlsx(path: Path) -> list[tuple[str, str]]:
    """Excel 2007+ (.xlsx) — extract cell text per sheet via openpyxl."""
    try:
        import openpyxl  # type: ignore
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        pairs: list[tuple[str, str]] = []
        for idx, sheet in enumerate(wb.worksheets, start=1):
            lines: list[str] = []
            for row in sheet.iter_rows(values_only=True):
                cells = [str(v) for v in row if v is not None and str(v).strip()]
                if cells:
                    lines.append("\t".join(cells))
            text = "\n".join(lines).strip()
            if text:
                pairs.append((f"sheet{idx:02d}", text))
        wb.close()
        return pairs
    except Exception as e:
        logger.warning("XLSX extraction failed for %s: %s", path.name, e)
        return []


def _chunks_xls(path: Path) -> list[tuple[str, str]]:
    """Excel 97-2003 (.xls) — extract cell text per sheet via xlrd."""
    try:
        import xlrd  # type: ignore
        wb = xlrd.open_workbook(str(path))
        pairs: list[tuple[str, str]] = []
        for idx in range(wb.nsheets):
            sheet = wb.sheet_by_index(idx)
            lines: list[str] = []
            for row_idx in range(sheet.nrows):
                cells = [
                    str(sheet.cell_value(row_idx, col))
                    for col in range(sheet.ncols)
                    if str(sheet.cell_value(row_idx, col)).strip()
                ]
                if cells:
                    lines.append("\t".join(cells))
            text = "\n".join(lines).strip()
            if text:
                pairs.append((f"sheet{idx + 1:02d}", text))
        return pairs
    except Exception as e:
        logger.warning("XLS extraction failed for %s: %s", path.name, e)
        return []


def _chunks_txt(path: Path) -> list[tuple[str, str]]:
    text = path.read_text(encoding="utf-8", errors="replace")
    lines = text.splitlines()

    sections: list[list[str]] = []
    current: list[str] = []

    for line in lines:
        # Split on # headings or blank lines
        if line.startswith("#") or (not line.strip() and current):
            if line.startswith("#"):
                if current:
                    sections.append(current)
                current = [line]
            else:
                # blank line — end current section
                sections.append(current)
                current = []
        else:
            if line.strip():
                current.append(line)

    if current:
        sections.append(current)

    if not sections:
        return [("sec01", text.strip())]

    pairs: list[tuple[str, str]] = []
    idx = 1
    for lines_block in sections:
        block_text = "\n".join(lines_block).strip()
        if not block_text:
            continue
        ref = f"sec{idx:02d}"
        if len(block_text) > MAX_TXT_CHARS:
            parts = _split_text(block_text, MAX_TXT_CHARS)
            pairs.extend(_apply_suffix(ref, parts))
        else:
            pairs.append((ref, block_text))
        idx += 1

    return pairs if pairs else [("sec01", text.strip())]


# ---------------------------------------------------------------------------
# run_extract — full pipeline step
# ---------------------------------------------------------------------------

def run_extract(
    conn: sqlite3.Connection,
    file_id: str,
    path: Path,
    lane: str,
    cfg_obj: Config | None = None,
) -> dict[str, int]:
    """
    Extract and persist chunks for one file.

    Returns:
        {"new": N, "updated": N, "unchanged": N, "skipped": N, "failed": N}

    Never raises.
    """
    stats: dict[str, int] = {
        "new": 0, "updated": 0, "unchanged": 0, "skipped": 0, "failed": 0
    }

    try:
        _cfg = cfg_obj or _module_cfg

        # --- Inference (always, regardless of lane) ---
        project_id, proj_confidence, signals = infer_project(path, cfg_obj=_cfg)
        typology = infer_typology(path, cfg_obj=_cfg)

        # Persist inference results to files row
        conn.execute(
            """UPDATE files
               SET project_id=?, project_confidence=?, typology=?, updated_at=datetime('now')
               WHERE file_id=?""",
            (project_id, proj_confidence, typology, file_id),
        )
        conn.commit()

        # Upsert projects table
        upsert_project(conn, {
            "project_id":     project_id,
            "typology_guess": typology,
            "signals":        json.dumps(signals),
        })

        # --- METADATA_ONLY: create a surrogate meta chunk ---
        if lane != "TEXT_EXTRACTABLE":
            surrogate = _build_surrogate(path, "metadata only")
            chunk_id = _make_chunk_id(file_id, "meta")
            c_hash = _content_hash(surrogate)
            existing = conn.execute(
                "SELECT content_hash FROM chunks WHERE chunk_id=?", (chunk_id,)
            ).fetchone()
            if existing:
                if existing["content_hash"] == c_hash:
                    stats["unchanged"] += 1
                else:
                    _upsert_chunk(conn, {
                        "chunk_id": chunk_id, "file_id": file_id,
                        "ref_value": "meta", "text": surrogate,
                        "token_estimate": _token_estimate(surrogate),
                        "content_hash": c_hash, "embedded": 0,
                    })
                    stats["updated"] += 1
            else:
                _upsert_chunk(conn, {
                    "chunk_id": chunk_id, "file_id": file_id,
                    "ref_value": "meta", "text": surrogate,
                    "token_estimate": _token_estimate(surrogate),
                    "content_hash": c_hash,
                })
                stats["new"] += 1
            set_file_status(conn, file_id, "EXTRACTED")
            log_event(conn, "EXTRACTED", detail="metadata only", file_id=file_id)
            return stats

        # --- TEXT_EXTRACTABLE: extract chunks ---
        chunk_pairs = extract_chunks(path)

        if not chunk_pairs:
            # Extraction returned nothing (e.g. scanned PDF with no text)
            set_file_status(conn, file_id, "FAILED", "EXTRACT_EMPTY",
                            "extract_chunks returned no chunks")
            log_event(conn, "EXTRACT_FAILED", detail="no chunks produced", file_id=file_id)
            stats["failed"] += 1
            return stats

        for ref, text in chunk_pairs:
            chunk_id = _make_chunk_id(file_id, ref)
            c_hash = _content_hash(text)
            tok_est = _token_estimate(text)

            existing = conn.execute(
                "SELECT content_hash FROM chunks WHERE chunk_id=?", (chunk_id,)
            ).fetchone()

            if existing:
                if existing["content_hash"] == c_hash:
                    stats["unchanged"] += 1
                else:
                    _upsert_chunk(conn, {
                        "chunk_id": chunk_id, "file_id": file_id,
                        "ref_value": ref, "text": text,
                        "token_estimate": tok_est,
                        "content_hash": c_hash, "embedded": 0,
                    })
                    stats["updated"] += 1
            else:
                _upsert_chunk(conn, {
                    "chunk_id": chunk_id, "file_id": file_id,
                    "ref_value": ref, "text": text,
                    "token_estimate": tok_est, "content_hash": c_hash,
                })
                stats["new"] += 1

        set_file_status(conn, file_id, "EXTRACTED")
        log_event(
            conn, "EXTRACTED",
            detail=f"{len(chunk_pairs)} chunks (new={stats['new']} updated={stats['updated']})",
            file_id=file_id,
        )

    except Exception as e:
        logger.exception("run_extract failed for file_id=%s path=%s: %s", file_id, path, e)
        try:
            set_file_status(conn, file_id, "FAILED", "EXTRACT_ERROR", str(e))
            log_event(conn, "EXTRACT_FAILED", detail=str(e), file_id=file_id)
        except Exception:
            pass
        stats["failed"] += 1

    return stats


# ---------------------------------------------------------------------------
# Legacy API (backward compat)
# ---------------------------------------------------------------------------

def extract(path: Path) -> tuple[str, str]:
    """
    Extract all text from a file as a single string.
    Returns (text, surrogate). Neither will be None.
    Errors are caught and logged; returns ("", surrogate) on failure.
    """
    ext = path.suffix.lower()
    if ext == ".pdf":
        return _extract_pdf(path)
    elif ext in (".docx", ".doc"):
        return _extract_docx(path)
    elif ext in (".pptx", ".ppt"):
        return _extract_pptx(path)
    elif ext in (".xlsx", ".xls"):
        pairs = _chunks_xlsx(path) if ext == ".xlsx" else _chunks_xls(path)
        text = "\n".join(t for _, t in pairs)
        return text, _build_surrogate(path, f"{len(pairs)} sheet(s)")
    elif ext in (".txt", ".md"):
        return _extract_text(path)
    else:
        return "", _build_surrogate(path, "no text extraction available")


def _extract_pdf(path: Path) -> tuple[str, str]:
    try:
        import pypdf  # type: ignore
        reader = pypdf.PdfReader(str(path), strict=False)
        pages = []
        for page in reader.pages:
            try:
                pages.append(page.extract_text() or "")
            except Exception:
                pages.append("")
        text = "\n".join(pages).strip()
        return text, _build_surrogate(path, f"{len(reader.pages)} pages")
    except Exception as e:
        logger.warning("PDF extraction failed for %s: %s", path.name, e)
        return "", _build_surrogate(path, "PDF extraction failed")


def _extract_docx(path: Path) -> tuple[str, str]:
    try:
        from docx import Document  # type: ignore
        doc = Document(str(path))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        text = "\n".join(paragraphs).strip()
        return text, _build_surrogate(path, f"{len(paragraphs)} paragraphs")
    except Exception as e:
        logger.warning("DOCX extraction failed for %s: %s", path.name, e)
        return "", _build_surrogate(path, "DOCX extraction failed")


def _extract_pptx(path: Path) -> tuple[str, str]:
    try:
        from pptx import Presentation  # type: ignore
        prs = Presentation(str(path))
        slides_text: list[str] = []
        for slide in prs.slides:
            parts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = para.text.strip()
                        if line:
                            parts.append(line)
            if parts:
                slides_text.append(" ".join(parts))
        text = "\n".join(slides_text).strip()
        return text, _build_surrogate(path, f"{len(prs.slides)} slides")
    except Exception as e:
        logger.warning("PPTX extraction failed for %s: %s", path.name, e)
        return "", _build_surrogate(path, "PPTX extraction failed")


def _extract_text(path: Path) -> tuple[str, str]:
    try:
        text = path.read_text(encoding="utf-8", errors="replace").strip()
        return text, _build_surrogate(path, "plain text")
    except Exception as e:
        logger.warning("Text extraction failed for %s: %s", path.name, e)
        return "", _build_surrogate(path, "text read failed")


def _build_surrogate(path: Path, note: str) -> str:
    parts = list(path.parts)
    context_parts = parts[-4:] if len(parts) >= 4 else parts
    breadcrumb = " > ".join(context_parts)
    return f"{path.stem} | {path.suffix.lstrip('.')} | {note} | {breadcrumb}"


def word_count(text: str) -> int:
    return len(text.split()) if text else 0


# ---------------------------------------------------------------------------
# Chunk 8 — Email extraction (.msg / .eml)
# ---------------------------------------------------------------------------

from dataclasses import dataclass, field as _dc_field

@dataclass
class EmailExtract:
    subject:        str
    sender:         str
    sender_email:   str
    recipients:     list[str]
    date:           str          # ISO datetime
    body_text:      str
    attachments:    list[str]    # filenames only
    thread_subject: str          # subject with Re:/Fwd: stripped
    direction:      str          # inbound | outbound | internal | unknown
    parties:        list[str]    # email domains involved


def _strip_reply_prefix(subject: str) -> str:
    """Remove Re:, Fwd:, RE:, FW: prefixes."""
    return re.sub(r'^(?:Re|Fwd?|Fw):\s*', '', subject, flags=re.IGNORECASE).strip()


def _detect_direction(
    sender_domain: str,
    recipient_domains: list[str],
    internal_domains: list[str],
) -> str:
    if not sender_domain:
        return "unknown"
    internal_set = {d.lower() for d in internal_domains}
    sender_internal = sender_domain.lower() in internal_set
    recipients_all_internal = all(d.lower() in internal_set for d in recipient_domains if d)
    if sender_internal and recipients_all_internal:
        return "internal"
    if sender_internal:
        return "outbound"
    return "inbound"


def _email_domain(email_addr: str) -> str:
    """Extract domain from email address."""
    if "@" in email_addr:
        return email_addr.split("@")[-1].strip().lower()
    return ""


def extract_email(file_path: str | Path) -> EmailExtract | None:
    """
    Extract metadata and body from .msg or .eml files.
    Returns None on unrecoverable failure.
    """
    path = Path(file_path)
    ext = path.suffix.lower()

    try:
        from config import cfg as _cfg
        internal_domains: list[str] = getattr(_cfg, "email_internal_domains", ["woha.com"])
    except Exception:
        internal_domains = ["woha.com"]

    try:
        if ext == ".msg":
            return _extract_msg(path, internal_domains)
        elif ext == ".eml":
            return _extract_eml(path, internal_domains)
        else:
            return None
    except Exception as e:
        logger.warning("extract_email failed for %s: %s", path.name, e)
        return None


def _extract_msg(path: Path, internal_domains: list[str]) -> EmailExtract:
    import extract_msg as _emsg  # type: ignore
    msg = _emsg.Message(str(path))

    subject = (msg.subject or "").strip()
    sender_name = (msg.sender or "").strip()
    sender_email_raw = (msg.senderEmail or "").strip()
    sender_domain = _email_domain(sender_email_raw)

    # Recipients
    to_list = [str(r).strip() for r in (msg.to or "").split(";") if str(r).strip()]
    cc_list = [str(r).strip() for r in (msg.cc or "").split(";") if str(r).strip()]
    recipients = [r for r in to_list + cc_list if r]

    recipient_domains = [_email_domain(r) for r in recipients]
    all_domains = list({d for d in [sender_domain] + recipient_domains if d})

    date_str = str(msg.date) if msg.date else ""
    body = (msg.body or "").strip()
    attachments = [str(a.longFilename or a.shortFilename or "") for a in (msg.attachments or [])]

    direction = _detect_direction(sender_domain, recipient_domains, internal_domains)

    return EmailExtract(
        subject=subject,
        sender=sender_name,
        sender_email=sender_email_raw,
        recipients=recipients,
        date=date_str,
        body_text=body[:50_000],  # cap body for indexing
        attachments=[a for a in attachments if a],
        thread_subject=_strip_reply_prefix(subject),
        direction=direction,
        parties=all_domains,
    )


def _extract_eml(path: Path, internal_domains: list[str]) -> EmailExtract:
    import email as _email_lib
    from email.header import decode_header as _decode_header

    def _decode_str(value) -> str:
        if value is None:
            return ""
        parts = _decode_header(str(value))
        result = []
        for part, charset in parts:
            if isinstance(part, bytes):
                result.append(part.decode(charset or "utf-8", errors="replace"))
            else:
                result.append(str(part))
        return " ".join(result).strip()

    raw = path.read_bytes()
    msg = _email_lib.message_from_bytes(raw)

    subject = _decode_str(msg.get("Subject", ""))
    from_raw = _decode_str(msg.get("From", ""))
    # Parse sender name and email
    import email.utils as _eu
    sender_name, sender_email_raw = _eu.parseaddr(from_raw)
    sender_domain = _email_domain(sender_email_raw)

    to_raw = _decode_str(msg.get("To", ""))
    cc_raw = _decode_str(msg.get("Cc", ""))
    recipients_raw = [r for r in to_raw.split(",") + cc_raw.split(",") if r.strip()]
    recipients = [r.strip() for r in recipients_raw]
    recipient_domains = [_email_domain(r) for r in recipients]
    all_domains = list({d for d in [sender_domain] + recipient_domains if d})

    date_str = _decode_str(msg.get("Date", ""))

    # Extract body text
    body = ""
    attachments: list[str] = []
    for part in msg.walk():
        ct = part.get_content_type()
        disposition = part.get("Content-Disposition", "")
        if "attachment" in disposition:
            fname = part.get_filename() or ""
            if fname:
                attachments.append(_decode_str(fname))
        elif ct == "text/plain" and not body:
            payload = part.get_payload(decode=True)
            if payload:
                charset = part.get_content_charset() or "utf-8"
                body = payload.decode(charset, errors="replace")

    direction = _detect_direction(sender_domain, recipient_domains, internal_domains)

    return EmailExtract(
        subject=subject,
        sender=sender_name,
        sender_email=sender_email_raw,
        recipients=recipients,
        date=date_str,
        body_text=body.strip()[:50_000],
        attachments=attachments,
        thread_subject=_strip_reply_prefix(subject),
        direction=direction,
        parties=all_domains,
    )


def get_email_content_tags(email_data: EmailExtract) -> list[str]:
    """Assign content_tags based on subject/body/parties for correspondence tagging."""
    tags: list[str] = []
    subject_lower = email_data.subject.lower()
    body_lower = email_data.body_text[:2000].lower()
    combined = subject_lower + " " + body_lower

    waiver_keywords = {"waiver", "deviation", "exemption", "relaxation", "dispensation"}
    if any(kw in combined for kw in waiver_keywords):
        tags.append("waiver_correspondence")

    authority_domains = {"ura.gov.sg", "bca.gov.sg", "scdf.gov.sg", "pub.gov.sg",
                         "nea.gov.sg", "lta.gov.sg", "hdb.gov.sg", "sla.gov.sg"}
    if any(domain in email_data.parties for domain in authority_domains):
        tags.append("authority_correspondence")

    submission_keywords = {"transmittal", "submission", "submitted"}
    if any(kw in combined for kw in submission_keywords):
        tags.append("submission_correspondence")

    meeting_keywords = {"minutes", "mom", "minute of meeting"}
    if any(kw in combined for kw in meeting_keywords):
        tags.append("meeting_correspondence")

    return tags


# ---------------------------------------------------------------------------
# Chunk 9 — Image classification
# ---------------------------------------------------------------------------

@dataclass
class ImageClass:
    image_type: str   # render | scanned_drawing | scanned_document | photo | diagram | unknown
    needs_ocr:  bool
    confidence: float


def classify_image(file_path: str | Path) -> ImageClass:
    """
    Classify an image by type (render vs scanned) based on path analysis.
    Returns ImageClass with needs_ocr flag.
    """
    path = Path(file_path)
    path_lower = str(path).lower().replace("\\", "/")
    name_lower = path.name.lower()

    # Render indicators (path-based, fast)
    render_path_signals = {"/viz/", "/renders/", "/cgi/", "/3d/", "/visualis", "/perspectives/",
                           "/enscape/", "/lumion/", "/twinmotion/", "\\viz\\", "\\renders\\"}
    if any(sig in path_lower for sig in render_path_signals):
        return ImageClass(image_type="render", needs_ocr=False, confidence=0.9)

    # Scan indicators
    scan_signals = {"/scan/", "/scanned/", "\\scan\\", "\\scanned\\"}
    if any(sig in path_lower for sig in scan_signals) or "scan" in name_lower:
        return ImageClass(image_type="scanned_document", needs_ocr=True, confidence=0.85)

    # Large file in Images folder → likely render
    try:
        size_mb = path.stat().st_size / (1024 * 1024)
        if size_mb > 5 and "/images/" in path_lower:
            return ImageClass(image_type="render", needs_ocr=False, confidence=0.75)
    except OSError:
        pass

    # Default: unknown, skip OCR (don't waste time)
    return ImageClass(image_type="unknown", needs_ocr=False, confidence=0.5)


def build_image_synthetic_description(file_path: str | Path, image_class: ImageClass) -> str:
    """
    Build a synthetic text description for renders/images so they're findable
    via semantic search without reading the image content.
    """
    path = Path(file_path)
    parts = path.parts

    # Extract useful tokens from path
    tokens: list[str] = []
    for part in parts[-6:]:
        cleaned = re.sub(r'[-_.]', ' ', part).strip()
        if cleaned and cleaned not in ("jpg", "jpeg", "png", "gif", "tif", "tiff"):
            tokens.append(cleaned)

    type_word = {
        "render": "render visualisation image",
        "photo": "photograph photo",
        "scanned_drawing": "scanned drawing",
        "scanned_document": "scanned document",
        "diagram": "diagram diagram",
    }.get(image_class.image_type, "image")

    desc = f"Project image {type_word} — {' '.join(tokens)}"
    return desc.strip()


# ---------------------------------------------------------------------------
# Chunk 12 — Content-based document classification
# ---------------------------------------------------------------------------

@dataclass
class ContentClass:
    doc_type:         str          # minutes | transmittal | authority_submission |
                                   # waiver_request | design_brief | specification | unknown
    confidence:       float
    detected_signals: list[str]


def classify_document_by_content(text_preview: str, filename: str = "") -> ContentClass:
    """
    Classify a document by reading the first 500 words.
    Fast pattern matching — no LLM required.

    Returns ContentClass with doc_type, confidence, and detected_signals.
    """
    # Normalise to first 500 words
    words = text_preview.split()[:500]
    preview = " ".join(words).lower()
    first_50 = " ".join(words[:50]).lower()
    first_20 = " ".join(words[:20]).lower()

    signals: list[str] = []

    # ── MINUTES ───────────────────────────────────────────────
    minutes_signals = 0
    if "minutes of meeting" in first_50 or "mom" in first_20:
        signals.append("mom_header"); minutes_signals += 2
    if re.search(r"present\s*:", preview) or "attendees:" in preview or "in attendance" in preview:
        signals.append("attendees_list"); minutes_signals += 1
    if re.search(r"action\s+by", preview) or "action items" in preview:
        signals.append("action_items"); minutes_signals += 1
    if "date of meeting" in preview or "venue:" in preview:
        signals.append("venue_date"); minutes_signals += 1
    if minutes_signals >= 2:
        conf = min(0.5 + minutes_signals * 0.1, 1.0)
        return ContentClass("minutes", conf, signals)

    # ── TRANSMITTAL ───────────────────────────────────────────
    trans_signals = 0
    if "transmittal" in first_20:
        signals.append("transmittal_header"); trans_signals += 3
    if "we transmit herewith" in preview or "please find enclosed" in preview:
        signals.append("transmittal_phrase"); trans_signals += 2
    if "document number" in preview and "revision" in preview:
        signals.append("doc_table"); trans_signals += 1
    if "transmitted by" in preview or "issued by" in preview:
        signals.append("transmitted_by"); trans_signals += 1
    if trans_signals >= 2:
        conf = min(0.5 + trans_signals * 0.1, 1.0)
        return ContentClass("transmittal", conf, signals)

    # ── AUTHORITY SUBMISSION ──────────────────────────────────
    auth_signals = 0
    for authority in ["ura", "bca", "scdf", "pub", "lta", "hdb"]:
        if f"submission to {authority}" in preview or f"to the {authority}" in first_50:
            signals.append(f"authority_{authority}"); auth_signals += 2; break
    if "building plan" in preview or "development application" in preview:
        signals.append("building_plan"); auth_signals += 1
    if re.search(r"ura/dc/|bca/|scdf/", preview):
        signals.append("authority_ref"); auth_signals += 2
    if auth_signals >= 2:
        return ContentClass("authority_submission", min(0.5 + auth_signals * 0.1, 1.0), signals)

    # ── WAIVER / DEVIATION REQUEST ────────────────────────────
    waiver_signals = 0
    if re.search(r"we (?:wish to apply|request|seek)\s+(?:a\s+)?(?:waiver|deviation|exemption)", preview):
        signals.append("waiver_request_phrase"); waiver_signals += 2
    if "in lieu of" in preview or "alternative solution" in preview:
        signals.append("in_lieu_of"); waiver_signals += 1
    if re.search(r"regulation\s+\d+", preview):
        signals.append("regulation_ref"); waiver_signals += 1
    if waiver_signals >= 2:
        return ContentClass("waiver_request", min(0.5 + waiver_signals * 0.1, 1.0), signals)

    # ── DESIGN BRIEF ──────────────────────────────────────────
    brief_signals = 0
    if "project brief" in first_20 or "design brief" in first_20:
        signals.append("brief_header"); brief_signals += 2
    if "scope of works" in preview or "project description" in preview:
        signals.append("scope_description"); brief_signals += 1
    if "design requirements" in preview:
        signals.append("design_requirements"); brief_signals += 1
    if brief_signals >= 2:
        return ContentClass("design_brief", min(0.5 + brief_signals * 0.1, 1.0), signals)

    # ── SPECIFICATION ─────────────────────────────────────────
    spec_signals = 0
    if "specification" in first_20:
        signals.append("spec_header"); spec_signals += 2
    if re.search(r"\b\d+\.\d+\.\d+\b", preview):  # clause numbering like 1.1.1
        signals.append("clause_numbering"); spec_signals += 1
    if "materials and workmanship" in preview:
        signals.append("workmanship"); spec_signals += 1
    if "contractor shall" in preview or "the contractor shall" in preview:
        signals.append("contractor_shall"); spec_signals += 1
    if spec_signals >= 2:
        return ContentClass("specification", min(0.5 + spec_signals * 0.1, 1.0), signals)

    return ContentClass("unknown", 0.3, signals)
